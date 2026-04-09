from pptx import Presentation
from pptx.util import Pt, Emu

prs = Presentation(r'poject seminar 1.pptx')

# Detailed analysis of slides 1-12 only
for slide_idx, slide in enumerate(prs.slides):
    if slide_idx >= 12:
        break
    print(f'\n=== SLIDE {slide_idx+1} [{slide.slide_layout.name}] ===')
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            if not para.text.strip():
                continue
            align = para.alignment
            for run in para.runs:
                f = run.font
                color = None
                try:
                    if f.color and f.color.type:
                        color = str(f.color.rgb)
                except:
                    pass
                sz = round(f.size / 12700, 1) if f.size else None
                txt = run.text[:80]
                print(f'  font="{f.name}" sz={sz}pt bold={f.bold} color={color} align={align} | "{txt}"')
