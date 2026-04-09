"""
Generate Seminar 2 PPT using the same template as Seminar 1.
Strategy: Open the template, copy its slide master/layouts into a new empty PPT,
then add our slides.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from copy import deepcopy
import os
import shutil
from lxml import etree

# Step 1: Create a clean presentation from the template
# We'll open the template, delete its slides via XML, then save
template_path = r'poject seminar 1.pptx'
output_path = r'project_seminar_2.pptx'
temp_path = r'_temp_clean.pptx'

# Copy the template
shutil.copy2(template_path, temp_path)

# Open and remove slides via XML manipulation
prs = Presentation(temp_path)
# Access the presentation XML
presentation_part = prs.part
nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
         'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
         'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}

# Get slide ID list element
sldIdLst = prs.presentation.sldIdLst
# Remove all slide references
for sldId in list(sldIdLst):
    rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    sldIdLst.remove(sldId)

prs.save(temp_path)

# Re-open the cleaned template
prs = Presentation(temp_path)
print(f'Slides after cleanup: {len(prs.slides)}')

# Get layouts
title_layout = None
title_only_layout = None
for layout in prs.slide_layouts:
    if layout.name == 'TITLE':
        title_layout = layout
    elif layout.name == 'TITLE_ONLY':
        title_only_layout = layout

if not title_layout:
    title_layout = prs.slide_layouts[0]
if not title_only_layout:
    title_only_layout = prs.slide_layouts[4] if len(prs.slide_layouts) > 4 else prs.slide_layouts[1]

print(f'Using title layout: {title_layout.name}')
print(f'Using title_only layout: {title_only_layout.name}')

# === STYLE CONSTANTS ===
FONT_TITLE = 'Fira Sans Extra Condensed SemiBold'
FONT_BODY = 'Fira Sans Extra Condensed'
FONT_MEDIUM = 'Fira Sans Extra Condensed Medium'

COLOR_ACCENT = RGBColor(0x40, 0x6E, 0x8E)
COLOR_ACCENT2 = RGBColor(0x2C, 0x3E, 0x50)
COLOR_GREEN = RGBColor(0x27, 0xAE, 0x60)
COLOR_LIGHT_GRAY = RGBColor(0x95, 0xA5, 0xA6)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_title(slide, text, font_size=None):
    """Set the title placeholder text"""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = text
            for para in shape.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                if font_size:
                    for run in para.runs:
                        run.font.size = Pt(font_size)
            return
    # Fallback: add text box as title
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size or 28)
    run.font.name = FONT_TITLE


def add_text(slide, text, left, top, width, height, font_size=13, bold=False,
             color=None, alignment=PP_ALIGN.LEFT, font_name=None):
    """Add a text box"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alignment
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.name = font_name or FONT_BODY
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
    return txBox


def add_bullets(slide, items, left, top, width, height=3.5, font_size=13,
                color=None, spacing=Pt(5)):
    """Add bullet points with smart bold-before-colon formatting"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = spacing
        
        if ': ' in item and not item.startswith('[') and not item.startswith('http'):
            bold_part, rest = item.split(': ', 1)
            r1 = p.add_run()
            r1.text = '• ' + bold_part + ': '
            r1.font.size = Pt(font_size)
            r1.font.name = FONT_BODY
            r1.font.bold = True
            if color:
                r1.font.color.rgb = color
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(font_size)
            r2.font.name = FONT_BODY
            if color:
                r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = '• ' + item
            r.font.size = Pt(font_size)
            r.font.name = FONT_BODY
            if color:
                r.font.color.rgb = color
    return txBox


def add_table(slide, data, left, top, width, row_height=0.35, font_size=10):
    """Add a table"""
    rows, cols = len(data), len(data[0])
    tbl = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                  Inches(width), Inches(row_height * rows))
    table = tbl.table
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size = Pt(font_size)
                    run.font.name = FONT_BODY
                    if ri == 0:
                        run.font.bold = True
                        run.font.color.rgb = COLOR_WHITE
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_ACCENT2
    return tbl


def add_img(slide, path, left, top, width=None, height=None):
    """Add an image if it exists"""
    if os.path.exists(path):
        kwargs = {}
        if width: kwargs['width'] = Inches(width)
        if height: kwargs['height'] = Inches(height)
        slide.shapes.add_picture(path, Inches(left), Inches(top), **kwargs)
        return True
    return False


# ========================================================================
# SLIDE 1: Title
# ========================================================================
s = prs.slides.add_slide(title_layout)
for shape in s.placeholders:
    idx = shape.placeholder_format.idx
    if idx == 0:  # Title
        shape.text = ""
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = "Distributed Feature Engineering for GNN-based "
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = "Fraud Detection System"
    elif idx == 1:  # Subtitle
        shape.text = ""
        tf = shape.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = "Group no. : 13"
        r.font.name = FONT_MEDIUM
        r.font.size = Pt(18)
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = "Guide : Prof. Dev Mukherjee"
        r2.font.name = FONT_MEDIUM
        r2.font.size = Pt(18)
        p3 = tf.add_paragraph()
        r3 = p3.add_run()
        r3.text = "Seminar 2 — Implementation & Results"
        r3.font.name = FONT_MEDIUM
        r3.font.size = Pt(14)

# ========================================================================
# SLIDE 2: Agenda
# ========================================================================
s = prs.slides.add_slide(title_only_layout)
set_title(s, "Agenda")
add_bullets(s, [
    "Problem Statement & Motivation",
    "Solution Overview — Why Graph Neural Networks?",
    "System Architecture & Design",
    "Heterogeneous Graph Construction",
    "GNN Architecture — HeteroRGCN",
    "DiFE Framework — Distributed Feature Engineering",
    "Feature Categories & Extractors",
    "Implementation Details & Tech Stack",
    "Training Results & Evaluation Metrics",
    "Key Challenges & Solutions",
    "Conclusion & Future Work",
], left=1.0, top=1.0, width=8, font_size=14, spacing=Pt(4))

# ========================================================================
# SLIDE 3: Problem Statement
# ========================================================================
s = prs.slides.add_slide(title_only_layout)
set_title(s, "Problem Statement")
add_bullets(s, [
    "Online payment fraud causes losses exceeding $40 billion annually worldwide",
    "Traditional ML models treat transactions independently — they miss relational patterns",
    "Fraudsters evolve techniques: card testing, identity theft, account takeover, fraud rings",
    "Need: A system that captures relationships between transactions, cards, addresses, and devices",
    "Challenge: Highly imbalanced data — only 3.5% of transactions are fraudulent",
], left=0.5, top=1.0, width=9, font_size=13)

add_text(s, "Key Insight: Fraud is inherently relational — fraudsters share cards, addresses,\ndevices, and emails across multiple transactions, forming detectable graph patterns.",
         left=0.5, top=3.7, width=9, height=0.8, font_size=12, bold=True, color=COLOR_ACCENT)

# ========================================================================
# SLIDE 4: Proposed Solution
# ========================================================================
s = prs.slides.add_slide(title_only_layout)
set_title(s, "Proposed Solution")
add_bullets(s, [
    "Graph Neural Network (GNN) on a heterogeneous transaction graph",
    "Model: Heterogeneous Relational Graph Convolutional Network (HeteroRGCN)",
    "Key Innovation: Distributed Feature Engineering (DiFE) framework",
    "Dataset: IEEE-CIS Fraud Detection (Kaggle) — 590K+ transactions",
    "Graph Scale: 726,345 nodes | 19,518,802 edges | 50+ node types",
], left=0.5, top=1.0, width=9, font_size=13)

add_text(s, "Why GNN over Traditional ML?", left=0.5, top=3.1, width=9, height=0.4,
         font_size=14, bold=True, color=COLOR_ACCENT)
add_bullets(s, [
    "Captures structural relationships between entities (cards, addresses, devices)",
    "Learns from both node features AND graph topology simultaneously",
    "Detects fraud rings and collusion patterns invisible to tabular classifiers",
], left=0.5, top=3.5, width=9, font_size=12, spacing=Pt(3))

# ========================================================================
# SLIDE 5: System Architecture
# ========================================================================
s = prs.slides.add_slide(title_only_layout)
set_title(s, "System Architecture")
arch = (
    "Raw Data (IEEE-CIS Transaction + Identity CSVs)\n"
    "         ↓\n"
    "Data Preprocessing & Feature Extraction\n"
    "         ↓\n"
    "Heterogeneous Graph Construction (PyG HeteroData)\n"
    "         ↓\n"
    "Distributed Feature Engineering (DiFE) — 7 Extractors in Parallel\n"
    "         ↓\n"
    "HeteroRGCN Model (3-layer RGCN with LeakyReLU)\n"
    "         ↓\n"
    "Binary Classification: Fraud / Not Fraud"
)
add_text(s, arch, left=1.5, top=0.9, width=7, height=4.0, font_size=14,
         alignment=PP_ALIGN.CENTER, font_name=FONT_MEDIUM)

# ========================================================================
# SLIDE 6: Heterogeneous Graph Construction
# ========================================================================
s = prs.slides.add_slide(title_only_layout)
set_title(s, "Heterogeneous Graph Construction")

add_text(s, "What is a Heterogeneous Graph?", left=0.5, top=0.8, width=5, height=0.4,
         font_size=14, bold=True, color=COLOR_ACCENT)
add_bullets(s, [
    "Multiple node types: transaction, card, address, email, device, product",
    "Multiple edge types: transaction→card, transaction→address, etc.",
    "Each transaction connects to its related entities via typed edges",
    "Allows capturing multi-hop relational patterns",
], left=0.5, top=1.2, width=5, font_size=12, spacing=Pt(3))

add_text(s, "Graph Statistics", left=5.5, top=0.8, width=4, height=0.4,
         font_size=14, bold=True, color=COLOR_ACCENT)
add_table(s, [
    ["Metric", "Value"],
    ["Nodes", "726,345"],
    ["Edges", "19,518,802"],
    ["Node Types", "50+"],
    ["Edge Types", "100+"],
    ["Fraud Rate", "3.5%"],
], left=5.5, top=1.2, width=4, font_size=11)

add_img(s, 'graph_intro.png', left=0.5, top=3.3, width=4.0, height=2.0)

# ========================================================================
# SLIDE 7: How Fraud Patterns Emerge
# ========================================================================
s = prs.slides.add_slide(title_only_layout)
set_title(s, "How Fraud Patterns Emerge in Graphs")
add_bullets(s, [
    "Card Testing Attack: Multiple small txns, same card, different merchants → rapid_transactions + small_amounts detected",
    "Fraud Ring: Multiple cards → same billing address → addr_card_diversity detects shared infrastructure",
    "Account Takeover: Same device, multiple card numbers → device_sharing feature fires",
    "Velocity Abuse: High frequency from single entity → temporal features flag anomalies",
    "GNN Advantage: Message passing propagates fraud signals across the graph — if neighbors are mostly fraud, flag the node",
], left=0.5, top=1.0, width=9, font_size=12, spacing=Pt(10))

# ========================================================================
# SLIDE 8: GNN Architecture
# ========================================================================
s = prs.slides.add_slide(title_only_layout)
set_title(s, "GNN Architecture — HeteroRGCN")

add_text(s, "Model Architecture", left=0.3, top=0.85, width=4.8, height=0.4,
         font_size=14, bold=True, color=COLOR_ACCENT)
add_text(s, (
    "Input: Node features (target) + Trainable embeddings (non-target)\n"
    "  ↓\n"
    "HeteroRGCN Layer 1 — Relation-specific Linear transforms\n"
    "  ↓ LeakyReLU\n"
    "HeteroRGCN Layer 2 — Message passing across relations\n"
    "  ↓ LeakyReLU\n"
    "HeteroRGCN Layer 3 — Final aggregation\n"
    "  ↓\n"
    "Output: Binary Classification (Fraud / Not Fraud)"
), left=0.3, top=1.2, width=4.8, height=3.0, font_size=11)

add_text(s, "Key Design Choices", left=5.3, top=0.85, width=4.5, height=0.4,
         font_size=14, bold=True, color=COLOR_ACCENT)
add_bullets(s, [
    "Sum aggregation across relation types",
    "Relation-specific weight matrices (W_r)",
    "Trainable embeddings for non-target nodes",
    "3 GNN layers for multi-hop aggregation",
], left=5.3, top=1.2, width=4.5, font_size=12, spacing=Pt(3))

add_table(s, [
    ["Parameter", "Value"],
    ["Hidden Size", "16"],
    ["Num Layers", "3"],
    ["Learning Rate", "0.01"],
    ["Embedding Size", "360"],
    ["Epochs", "1000"],
], left=5.5, top=3.0, width=3.5, font_size=10)

# ========================================================================
# SLIDE 9: DiFE Framework Overview
# ========================================================================
s = prs.slides.add_slide(title_only_layout)
set_title(s, "DiFE Framework — Distributed Feature Engineering")

add_text(s, "Core Contribution: A modular, parallelizable feature engineering framework",
         left=0.5, top=0.85, width=9, height=0.4, font_size=13, bold=True, color=COLOR_ACCENT)
add_bullets(s, [
    "Modular Design: Each feature extractor is an independent, pluggable module",
    "Parallel Processing: Extractors run concurrently using Python multiprocessing",
    "Registry System: Feature extractors register with a centralized FeatureRegistry",
    "Caching: Extracted features cached to disk with hash-based invalidation",
    "Scalable: Supports 1 to N workers, distributing extraction across CPU cores",
    "Extensible: New extractors added by subclassing FeatureExtractor base class",
    "Integrated: Connects to training pipeline via --use-distributed-features flag",
], left=0.5, top=1.3, width=9, font_size=12, spacing=Pt(5))

add_text(s, "DiFE Flow: Transaction Data → FeatureRegistry → DistributedEngine (N workers) → Feature Combiner → Enhanced Matrix",
         left=0.3, top=4.2, width=9.4, height=0.5, font_size=12, bold=True,
         color=COLOR_ACCENT2, alignment=PP_ALIGN.CENTER)

# ========================================================================
# SLIDE 10: Feature Extractors Table
# ========================================================================
s = prs.slides.add_slide(title_only_layout)
set_title(s, "Feature Extractors — 7 Specialized Modules")

add_table(s, [
    ["#", "Extractor", "Type", "What It Captures"],
    ["1", "GraphNeighborAggregator", "Graph", "Neighbor statistics (mean/std/max/min/sum/count)"],
    ["2", "TemporalFeatureExtractor", "Temporal", "Time-of-day, day-of-week, velocity, hour concentration"],
    ["3", "StatisticalFeatureExtractor", "Statistical", "Skewness, kurtosis, entropy, distribution shape"],
    ["4", "RiskScoreExtractor", "Risk", "Amount anomalies, entity frequency, z-scores"],
    ["5", "GraphCentralityExtractor", "Graph", "Degree centrality, normalized centrality measures"],
    ["6", "PatternMatchingExtractor", "Risk", "Card testing, device sharing, round amounts, rapid txns"],
    ["7", "CrossFeatureExtractor", "Statistical", "Interactions: Amount×Product, Amount×Card, Dist×Addr"],
], left=0.3, top=1.0, width=9.4, font_size=10, row_height=0.38)

add_text(s, "Total: 470+ engineered features extracted per transaction",
         left=0.3, top=4.3, width=9.4, height=0.4, font_size=14, bold=True,
         color=COLOR_ACCENT, alignment=PP_ALIGN.CENTER)

# ========================================================================
# SLIDE 11: Feature Engineering Example
# ========================================================================
s = prs.slides.add_slide(title_only_layout)
set_title(s, "Feature Engineering — Detecting Card Testing")

add_text(s, "Example: How DiFE features detect a card testing attack",
         left=0.5, top=0.85, width=9, height=0.3, font_size=13, bold=True, color=COLOR_ACCENT)
add_bullets(s, [
    "Step 1 — rapid_transactions = 1: Multiple transactions in < 60 seconds",
    "Step 2 — small_amounts = 1: Transaction amount < $2 (testing if card works)",
    "Step 3 — round_amounts = 1: Perfect dollar amounts (no cents)",
    "Step 4 — card_velocity = HIGH: Same card used 10+ times in an hour",
    "Step 5 — neighbor_count = HIGH: Graph shows card connected to many merchants",
    "Step 6 — degree_centrality = HIGH: Card is a hub node in the transaction graph",
], left=0.5, top=1.2, width=9, font_size=12, spacing=Pt(6))

add_text(s, "Result: GNN combines graph topology + DiFE features → flags card testing with high confidence",
         left=0.5, top=4.0, width=9, height=0.5, font_size=13, bold=True, color=COLOR_GREEN)

# ========================================================================
# SLIDE 12: Implementation Details
# ========================================================================
s = prs.slides.add_slide(title_only_layout)
set_title(s, "Implementation Details")

add_text(s, "Technology Stack", left=0.3, top=0.85, width=4.5, height=0.3,
         font_size=14, bold=True, color=COLOR_ACCENT)
add_bullets(s, [
    "Python 3.8+",
    "PyTorch — Deep learning",
    "PyTorch Geometric (PyG) — GNN",
    "DGL — Graph processing",
    "Pandas & NumPy — Data wrangling",
    "scikit-learn — Metrics",
    "Multiprocessing — Parallelism",
    "Matplotlib — Visualization",
], left=0.3, top=1.2, width=4.5, font_size=11, spacing=Pt(2))

add_text(s, "Project Structure", left=5.0, top=0.85, width=4.8, height=0.3,
         font_size=14, bold=True, color=COLOR_ACCENT)
add_text(s, (
    "graph-fraud-detection/\n"
    "├── train.py\n"
    "├── distributed_feature_extraction.py\n"
    "├── gnn/\n"
    "│   ├── pytorch_model.py\n"
    "│   ├── feature_engineering.py\n"
    "│   ├── advanced_features.py\n"
    "│   ├── distributed_feature_pipeline.py\n"
    "│   ├── graph_utils.py\n"
    "│   ├── data.py\n"
    "│   └── utils.py\n"
    "├── ieee-data/\n"
    "├── output/\n"
    "└── model/"
), left=5.0, top=1.2, width=4.8, height=3.5, font_size=10)

# ========================================================================
# SLIDE 13: Data Preprocessing Pipeline
# ========================================================================
s = prs.slides.add_slide(title_only_layout)
set_title(s, "Data Preprocessing Pipeline")
add_bullets(s, [
    "Load IEEE-CIS transaction CSV — 590,540 rows × 394 columns",
    "Load identity CSV — device, browser, OS, IP information",
    "Join on TransactionID — merge transaction and identity data",
    "Extract edge lists — map relationships (card, address, email, device, product)",
    "Build heterogeneous graph — PyG HeteroData with typed nodes and edges",
    "Compute node features — numeric features normalized (mean=0, std=1)",
    "Train/test split — temporal split preserving fraud distribution",
    "Apply DiFE extractors — enhance with 470+ engineered features",
], left=0.5, top=1.0, width=9, font_size=13, spacing=Pt(7))

# ========================================================================
# SLIDE 14: Training Results — Loss & F1
# ========================================================================
s = prs.slides.add_slide(title_only_layout)
set_title(s, "Training Results — Loss & F1 Progression")

if not add_img(s, 'output/loss.jpg', left=0.3, top=1.0, width=4.4, height=3.0):
    add_text(s, "[Loss Curve — output/loss.jpg]", left=0.3, top=2.0, width=4.4, height=0.5,
             font_size=14, alignment=PP_ALIGN.CENTER, color=COLOR_LIGHT_GRAY)

if not add_img(s, 'output/f1.jpg', left=5.0, top=1.0, width=4.4, height=3.0):
    add_text(s, "[F1 Curve — output/f1.jpg]", left=5.0, top=2.0, width=4.4, height=0.5,
             font_size=14, alignment=PP_ALIGN.CENTER, color=COLOR_LIGHT_GRAY)

add_text(s, "Loss: 0.65 → 0.08  |  F1: 0.0 → peak 0.56  |  Training: ~31 sec/epoch × 1000 epochs ≈ 8.5 hours",
         left=0.3, top=4.2, width=9.4, height=0.5, font_size=11, bold=True, alignment=PP_ALIGN.CENTER)

# ========================================================================
# SLIDE 15: Results — ROC & PR
# ========================================================================
s = prs.slides.add_slide(title_only_layout)
set_title(s, "Results — ROC & Precision-Recall Curves")

if not add_img(s, 'output/roc_curve.png', left=0.3, top=1.0, width=4.4, height=3.2):
    add_text(s, "[ROC Curve — AUC = 0.92]", left=0.3, top=2.0, width=4.4, height=0.5,
             font_size=14, alignment=PP_ALIGN.CENTER, color=COLOR_LIGHT_GRAY)

if not add_img(s, 'output/pr_curve.png', left=5.0, top=1.0, width=4.4, height=3.2):
    add_text(s, "[PR Curve — AP = 0.59]", left=5.0, top=2.0, width=4.4, height=0.5,
             font_size=14, alignment=PP_ALIGN.CENTER, color=COLOR_LIGHT_GRAY)

add_text(s, "ROC AUC = 0.92 (strong discriminative ability)  |  PR AUC = 0.59 (effective on imbalanced data)",
         left=0.3, top=4.3, width=9.4, height=0.5, font_size=12, bold=True,
         color=COLOR_ACCENT, alignment=PP_ALIGN.CENTER)

# ========================================================================
# SLIDE 16: Performance Metrics
# ========================================================================
s = prs.slides.add_slide(title_only_layout)
set_title(s, "Performance Metrics Summary")

add_table(s, [
    ["Metric", "Value"],
    ["ROC AUC", "0.92"],
    ["Precision", "0.86"],
    ["Recall", "0.35"],
    ["F1 Score", "0.56 (best)"],
    ["PR AUC / AP", "0.59"],
    ["Accuracy", "97.6%"],
], left=0.5, top=1.0, width=4, font_size=12, row_height=0.4)

add_text(s, "Confusion Matrix", left=5.0, top=0.85, width=4.5, height=0.3,
         font_size=14, bold=True, color=COLOR_ACCENT)
add_table(s, [
    ["", "Actual Positive", "Actual Negative"],
    ["Predicted Positive", "1,435 (TP)", "240 (FP)"],
    ["Predicted Negative", "2,629 (FN)", "113,804 (TN)"],
], left=5.0, top=1.2, width=4.5, font_size=11, row_height=0.4)

add_text(s, (
    "High Precision (0.86): Only 240 false positives out of 1,675 flagged transactions\n"
    "Recall Trade-off (0.35): 2,629 fraud cases missed — room for improvement\n"
    "Design Choice: Prioritized precision to avoid disrupting legitimate users"
), left=0.5, top=3.7, width=9, height=1.0, font_size=12)

# ========================================================================
# SLIDE 17: DiFE Impact
# ========================================================================
s = prs.slides.add_slide(title_only_layout)
set_title(s, "DiFE Impact — Baseline vs Enhanced Features")

add_table(s, [
    ["Metric", "Baseline (No DiFE)", "With DiFE", "Improvement"],
    ["ROC AUC", "0.92", "TBD", "—"],
    ["Precision", "0.86", "TBD", "—"],
    ["Recall", "0.35", "TBD", "—"],
    ["F1 Score", "0.56", "TBD", "—"],
    ["PR AUC", "0.59", "TBD", "—"],
], left=1.0, top=1.2, width=8, font_size=13, row_height=0.42)

add_text(s, (
    "DiFE adds 470+ engineered features to the base feature set\n"
    "Expected: Improved recall and F1 while maintaining high precision\n"
    "Hypothesis: Graph-based and risk features capture patterns missed by raw features"
), left=0.5, top=4.0, width=9, height=1.0, font_size=12, color=COLOR_ACCENT)

# ========================================================================
# SLIDE 18: Challenges & Solutions
# ========================================================================
s = prs.slides.add_slide(title_only_layout)
set_title(s, "Key Challenges & Solutions")

add_table(s, [
    ["Challenge", "Solution"],
    ["Highly imbalanced data (3.5% fraud)", "Weighted cross-entropy loss, precision-focused evaluation"],
    ["Large graph (19.5M edges)", "Full-graph training with PyG, sparse operations"],
    ["Training instability (F1 collapse epoch 650+)", "Early stopping based on validation F1"],
    ["Feature dimensionality (470+ features)", "Feature selection, normalization, caching"],
    ["Scalability for large datasets", "DiFE: distributed parallel feature extraction"],
    ["Multiple entity types & relations", "Heterogeneous graph with relation-specific transforms"],
], left=0.3, top=1.0, width=9.4, font_size=10, row_height=0.45)

# ========================================================================
# SLIDE 19: Scalability
# ========================================================================
s = prs.slides.add_slide(title_only_layout)
set_title(s, "Distributed Processing & Scalability")

add_bullets(s, [
    "Each feature extractor runs independently — no inter-extractor dependencies",
    "Python multiprocessing distributes extractors across available CPU cores",
    "Features cached to disk with hash-based invalidation — avoids redundant computation",
    "Batch processing support for streaming / large-scale inference scenarios",
    "Framework extensible to distributed clusters (Spark, Dask, Ray)",
], left=0.5, top=1.0, width=9, font_size=13, spacing=Pt(8))

add_table(s, [
    ["Workers", "Feature Extraction Time", "Speedup"],
    ["1 (Sequential)", "Baseline", "1.0x"],
    ["2", "~50% reduction", "~2.0x"],
    ["4", "~75% reduction", "~3.5x"],
    ["8", "~87% reduction", "~6.5x"],
], left=2.0, top=3.5, width=6, font_size=11, row_height=0.35)

# ========================================================================
# SLIDE 20: Limitations
# ========================================================================
s = prs.slides.add_slide(title_only_layout)
set_title(s, "Limitations")
add_bullets(s, [
    "Recall at 0.35: Model misses ~65% of actual fraud cases — needs improvement",
    "Full-graph training: Entire graph must fit in memory — limits scalability",
    "No real-time inference: Current pipeline requires offline batch processing",
    "Static graph: Model must be retrained for new transactions (no incremental updates)",
    "Feature cost: 470+ features increase computation time and memory usage",
    "No interpretability: Black-box predictions — cannot explain why a transaction is flagged",
    "Single dataset: Results only validated on IEEE-CIS data",
], left=0.5, top=1.0, width=9, font_size=13, spacing=Pt(6))

# ========================================================================
# SLIDE 21: Future Work
# ========================================================================
s = prs.slides.add_slide(title_only_layout)
set_title(s, "Future Work")
add_bullets(s, [
    "Graph Attention Networks (GAT): Replace RGCN with attention-based aggregation for interpretability",
    "Mini-batch Training: NeighborLoader (PyG) for scalable training on large graphs",
    "Real-time Streaming: Incremental graph updates for online fraud detection",
    "Explainability: GNNExplainer / SHAP for human-readable flagging reasons",
    "Production API: FastAPI service for real-time fraud scoring",
    "Temporal Graph Networks: Dynamic GNNs capturing time-evolving fraud patterns",
    "Transfer Learning: Adapt model to insurance, healthcare, banking fraud domains",
], left=0.5, top=1.0, width=9, font_size=12, spacing=Pt(7))

# ========================================================================
# SLIDE 22: Conclusion
# ========================================================================
s = prs.slides.add_slide(title_only_layout)
set_title(s, "Conclusion")
add_bullets(s, [
    "Built a complete graph-based fraud detection system using HeteroRGCN",
    "Designed DiFE — distributed feature framework with 7 extractors, 470+ features",
    "Achieved ROC AUC of 0.92 and Precision of 0.86 on IEEE-CIS dataset",
    "Graph structure captures relational fraud patterns that tabular ML misses",
    "DiFE enables scalable, parallel feature extraction for deployment",
    "Identified improvement areas: recall enhancement, real-time inference, interpretability",
], left=0.5, top=1.0, width=9, font_size=14, spacing=Pt(10))

add_text(s, "GNN + DiFE = Scalable, accurate, and extensible fraud detection",
         left=0.5, top=4.2, width=9, height=0.5, font_size=16, bold=True,
         color=COLOR_ACCENT, alignment=PP_ALIGN.CENTER)

# ========================================================================
# SLIDE 23: References
# ========================================================================
s = prs.slides.add_slide(title_only_layout)
set_title(s, "References")
add_bullets(s, [
    "[1] IEEE-CIS Fraud Detection — Kaggle Competition Dataset",
    "[2] Kipf & Welling, \"Semi-Supervised Classification with GCNs\" (ICLR 2017)",
    "[3] Schlichtkrull et al., \"Modeling Relational Data with GCNs\" (ESWC 2018)",
    "[4] Wang et al., \"Heterogeneous Graph Attention Network\" (WWW 2019)",
    "[5] Liu et al., \"GeniePath: GNNs with Adaptive Receptive Paths\" (AAAI 2019)",
    "[6] AWS SageMaker — Graph-Based Fraud Detection (GitHub)",
    "[7] Deep Graph Library (DGL) — Official Documentation",
    "[8] PyTorch Geometric (PyG) — Official Documentation",
    "[9] Hamilton et al., \"Inductive Representation Learning on Large Graphs\" (NeurIPS 2017)",
    "[10] Chami et al., \"ML on Graphs: Model and Taxonomy\" (JMLR 2022)",
], left=0.3, top=0.9, width=9.4, font_size=11, spacing=Pt(3))

# ========================================================================
# SLIDE 24: Thank You
# ========================================================================
s = prs.slides.add_slide(title_only_layout)
set_title(s, "Thank You")
add_text(s, "Questions?", left=0.5, top=2.0, width=9, height=0.8,
         font_size=28, alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE, color=COLOR_ACCENT)
add_text(s, "Group 13 | Guide: Prof. Dev Mukherjee\nSeminar 2 — Implementation & Results",
         left=0.5, top=3.2, width=9, height=0.8, font_size=14,
         alignment=PP_ALIGN.CENTER, font_name=FONT_MEDIUM)


# ========================================================================
# SAVE
# ========================================================================
prs.save(output_path)
print(f'\n✅ Presentation saved: {output_path}')
print(f'   Total slides: {len(prs.slides)}')

# Cleanup
try:
    os.remove(temp_path)
except:
    pass

print('Done!')
